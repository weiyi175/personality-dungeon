from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/personality_dungeon_onboarding_full_deck.pptx")


class Theme:
    bg = RGBColor(248, 250, 252)
    title = RGBColor(15, 23, 42)
    body = RGBColor(30, 41, 59)
    accent = RGBColor(14, 116, 144)
    accent2 = RGBColor(4, 120, 87)
    warn = RGBColor(180, 83, 9)
    bad = RGBColor(190, 24, 93)
    good = RGBColor(5, 150, 105)
    box = RGBColor(226, 232, 240)


HIGHLIGHT_COLORS = {
    "info": RGBColor(219, 234, 254),
    "success": RGBColor(220, 252, 231),
    "warning": RGBColor(254, 240, 138),
    "danger": RGBColor(254, 226, 226),
}


def style_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = Theme.bg


def add_title_box(slide, title: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = Theme.title

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.18), Inches(11.4), Inches(0.7))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Microsoft JhengHei"
        sp.font.size = Pt(17)
        sp.font.color.rgb = Theme.accent


def add_bullets(slide, lines: list[str], x=0.8, y=1.9, w=6.2, h=4.8, size=21):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(size)
        p.font.color.rgb = Theme.body
        p.space_after = Pt(10)


def add_rounded_box(slide, x, y, w, h, text, fill_rgb, font_size=16, bold=True):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.color.rgb = RGBColor(148, 163, 184)
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = Theme.title
    return shp


def add_highlight_box(slide, x, y, w, h, text, kind="info", font_size=16):
    fill_rgb = HIGHLIGHT_COLORS.get(kind, HIGHLIGHT_COLORS["info"])
    return add_rounded_box(slide, x, y, w, h, text, fill_rgb, font_size=font_size, bold=True)


def add_arrow(slide, x, y, w, h, color=RGBColor(100, 116, 139)):
    arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.color.rgb = color
    return arr


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = RGBColor(100, 116, 139)
    line.line.width = Pt(2)
    return line


def add_footer(slide, idx: int):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(11.5), Inches(0.3))
    tf = b.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Personality Dungeon | Onboarding Deck | p.{idx}"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def add_auto_paged_bullets(prs, slide_no, title, subtitle, lines, x=0.9, y=1.9, w=11.2, h=4.8, size=18, max_lines=4):
    chunks = [lines[i:i + max_lines] for i in range(0, len(lines), max_lines)]
    for idx, chunk in enumerate(chunks):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        style_slide_bg(s)
        suffix = "（續）" if idx > 0 else ""
        add_title_box(s, f"{title}{suffix}", subtitle)
        add_bullets(s, chunk, x=x, y=y, w=w, h=h, size=size)
        add_footer(s, slide_no)
        slide_no += 1
    return slide_no


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_no = 1

    # 1. Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "個性地下城：人格 × 演化 × 遊戲系統", "Dynamic AI Ecology System | 新組員導覽完整版")
    add_rounded_box(s, 0.9, 2.0, 3.3, 1.5, "玩家\nIndividual Evolver", RGBColor(219, 234, 254))
    add_rounded_box(s, 4.95, 2.0, 3.3, 1.5, "地下城\nPersonality Testbed", RGBColor(209, 250, 229))
    add_rounded_box(s, 8.95, 2.0, 3.3, 1.5, "小火龍\nMeta Controller", RGBColor(254, 240, 138))
    add_arrow(s, 4.25, 2.45, 0.65, 0.45)
    add_arrow(s, 8.25, 2.45, 0.65, 0.45)
    add_bullets(s, ["作者：研究小組", "日期：2026-04-17", "關鍵主題：Level 2 plateau 與 Level 3 emergence"], x=0.95, y=4.2, w=11, h=1.8, size=18)
    add_footer(s, slide_no)
    slide_no += 1

    # 2. One sentence
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "一句話理解", "玩家在學習，但世界也在學習")
    add_bullets(
        s,
        [
            "玩家 = 個體演化單元（strategy selection + utility update）",
            "地下城 = 人格測試場（payoff feedback）",
            "小火龍 = 世界壓力生成器（反收斂干預）",
        ],
        x=0.8,
        y=2.0,
        w=6.4,
        h=3.8,
        size=20,
    )
    add_rounded_box(s, 8.2, 2.0, 4.2, 1.1, "玩家", RGBColor(219, 234, 254), font_size=20)
    add_rounded_box(s, 8.2, 3.4, 4.2, 1.1, "地下城", RGBColor(209, 250, 229), font_size=20)
    add_rounded_box(s, 8.2, 4.8, 4.2, 1.1, "小火龍", RGBColor(254, 240, 138), font_size=20)
    add_footer(s, slide_no)
    slide_no += 1

    # 3. Gameplay loop
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "玩家如何遊玩：核心 Loop", "選擇策略 → 挑戰 → 回饋 → 人格更新")
    labels = ["1. 選擇策略", "2. 挑戰地下城", "3. 成功/失敗", "4. 權重更新", "5. 進入下一輪"]
    x = 0.8
    for i, lb in enumerate(labels):
        add_rounded_box(s, x, 2.8, 2.2, 1.2, lb, RGBColor(226, 232, 240), font_size=16)
        if i < len(labels) - 1:
            add_arrow(s, x + 2.25, 3.15, 0.55, 0.45)
        x += 2.55
    add_bullets(s, ["每輪都會留下可分析的 timeseries 資料：p_*、w_*、avg_reward"], x=0.9, y=4.8, w=11.5, h=1.2, size=17)
    add_footer(s, slide_no)
    slide_no += 1

    # 4. Challenge mechanism
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "挑戰機制：玩家 A 挑戰玩家 B 的地下城")
    add_rounded_box(s, 1.0, 2.4, 2.5, 1.0, "玩家 A", RGBColor(219, 234, 254), font_size=18)
    add_arrow(s, 3.7, 2.65, 1.0, 0.45)
    add_rounded_box(s, 5.0, 2.2, 3.2, 1.4, "玩家 B 地下城", RGBColor(209, 250, 229), font_size=18)
    add_arrow(s, 8.35, 2.65, 1.0, 0.45)
    add_highlight_box(s, 9.55, 1.8, 2.7, 1.2, "成功\n樣本+金幣", kind="success", font_size=16)
    add_highlight_box(s, 9.55, 3.3, 2.7, 1.2, "失敗\n三層損失", kind="danger", font_size=16)
    add_bullets(
        s,
        [
            "失敗 = 三層損失：經濟層（金幣消耗）、學習層（樣本外流）、動力學層（權重劣化）",
            "這三層分別對應：經濟模型 / 資訊模型 / replicator 更新",
        ],
        x=0.95,
        y=5.0,
        w=11.2,
        h=1.4,
        size=16,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 4.1 Math-to-game mapping (critical bridge)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "數學如何對應玩家行為", "把 U = A x 直接對回 gameplay")
    add_highlight_box(s, 0.9, 2.0, 3.0, 1.0, "x(t)\n玩家策略比例", kind="info", font_size=17)
    add_arrow(s, 4.0, 2.3, 0.7, 0.4)
    add_highlight_box(s, 4.9, 2.0, 3.0, 1.0, "U = A x\n地下城 reward", kind="warning", font_size=17)
    add_arrow(s, 8.0, 2.3, 0.7, 0.4)
    add_highlight_box(s, 8.9, 2.0, 3.4, 1.0, "reward 高\n該策略玩家變多", kind="success", font_size=17)
    add_arrow(s, 6.2, 3.35, 0.7, 0.4)
    add_highlight_box(s, 4.9, 4.0, 3.4, 1.1, "分布改變\n下一輪 reward 改變", kind="info", font_size=16)
    add_bullets(
        s,
        [
            "核心直覺：玩家在玩遊戲，同時也在解一個延遲動態方程。",
            "這就是遊戲層、數學層、演化層的共同語言。",
        ],
        x=0.95,
        y=5.4,
        w=11.2,
        h=1.2,
        size=17,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 5. Dungeon nature
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "地下城的本質：人格外化測試場")
    add_rounded_box(s, 1.2, 2.1, 3.1, 1.2, "人格向量\n+ 偏好結構", RGBColor(219, 234, 254), font_size=17)
    add_arrow(s, 4.55, 2.45, 0.95, 0.45)
    add_rounded_box(s, 5.7, 2.1, 3.1, 1.2, "生成地下城規則\n事件/風險/回饋", RGBColor(254, 240, 138), font_size=16)
    add_arrow(s, 9.05, 2.45, 0.95, 0.45)
    add_rounded_box(s, 10.2, 2.1, 2.0, 1.2, "玩家體驗", RGBColor(209, 250, 229), font_size=17)
    add_bullets(
        s,
        [
            "不是傳統關卡腳本，而是『人格在環境中的可計算投影』",
            "可作為研究：策略偏好如何塑造可玩與可學習的世界",
        ],
        x=1.0,
        y=4.2,
        w=11.3,
        h=1.8,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 6. Meta GAI
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "小火龍（Meta GAI）：維持多樣性與壓力")
    add_bullets(
        s,
        [
            "全局分析：人格分布、勝率地圖、策略流行度、多樣性指標",
            "變異生成：事件模板、高難度地下城、風險模型",
            "壓力輸出：反收斂干預，避免單一策略壟斷",
            "數學化理解：小火龍 ≈ 動態修改 payoff structure 的外部控制器",
        ],
        x=0.8,
        y=1.9,
        w=7.2,
        h=3.8,
        size=17,
    )
    center = add_rounded_box(s, 8.4, 2.4, 3.4, 1.2, "小火龍控制器", RGBColor(254, 240, 138), font_size=18)
    _ = center
    add_rounded_box(s, 8.0, 4.2, 1.8, 0.9, "玩家A", RGBColor(219, 234, 254), font_size=14)
    add_rounded_box(s, 10.0, 4.2, 1.8, 0.9, "玩家B", RGBColor(219, 234, 254), font_size=14)
    add_rounded_box(s, 9.0, 5.3, 1.8, 0.9, "玩家C", RGBColor(219, 234, 254), font_size=14)
    add_connector(s, 10.1, 3.6, 8.9, 4.2)
    add_connector(s, 10.1, 3.6, 10.9, 4.2)
    add_connector(s, 10.1, 3.6, 9.9, 5.3)
    add_footer(s, slide_no)
    slide_no += 1

    # 6.1 Legend slide (low-cost high-value)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "圖例（Legend）", "統一顏色語意，降低理解成本")
    add_highlight_box(s, 1.0, 2.0, 3.0, 1.0, "藍：玩家", kind="info", font_size=18)
    add_highlight_box(s, 1.0, 3.3, 3.0, 1.0, "綠：地下城/正向", kind="success", font_size=18)
    add_highlight_box(s, 1.0, 4.6, 3.0, 1.0, "黃：小火龍/控制", kind="warning", font_size=18)
    add_highlight_box(s, 4.5, 2.0, 3.0, 1.0, "紅：風險/失敗", kind="danger", font_size=18)
    add_highlight_box(s, 4.5, 3.3, 3.0, 1.0, "灰：中性流程元件", kind="info", font_size=18)
    add_bullets(
        s,
        [
            "報告期間請盡量保持此顏色語意一致，避免跨頁切換時的認知成本。",
        ],
        x=0.95,
        y=5.2,
        w=11.2,
        h=1.0,
        size=17,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 7. State model
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "數學模型 I：狀態表示")
    add_bullets(
        s,
        [
            "狀態向量：x(t) = (x_A, x_D, x_B)",
            "約束：x_A + x_D + x_B = 1",
            "策略：aggressive / defensive / balanced",
            "資料來源：每輪抽樣比例 p_*（timeseries CSV）",
        ],
        x=0.9,
        y=1.9,
        w=6.3,
        h=3.8,
        size=19,
    )
    tri = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(8.0), Inches(2.1), Inches(4.3), Inches(3.8))
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(224, 242, 254)
    tri.line.color.rgb = RGBColor(14, 116, 144)
    for txt, tx, ty in [("A", 8.15, 5.5), ("D", 11.85, 5.5), ("B", 9.95, 2.2)]:
        t = s.shapes.add_textbox(Inches(tx), Inches(ty), Inches(0.4), Inches(0.4))
        p = t.text_frame.paragraphs[0]
        p.text = txt
        p.font.name = "Calibri"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = Theme.title
    add_footer(s, slide_no)
    slide_no += 1

    # 8. Payoff matrix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "數學模型 II：Payoff 機制（矩陣）")
    add_bullets(s, ["核心方程：U = A x", "其中 A 為循環對抗矩陣（RPS-like）"], x=0.9, y=1.8, w=5.5, h=1.5, size=20)
    matrix_text = "A = [ 0   a  -b ]\n    [ -b  0   a ]\n    [ a  -b  0 ]"
    m = s.shapes.add_textbox(Inches(0.95), Inches(2.9), Inches(5.7), Inches(2.2))
    mp = m.text_frame.paragraphs[0]
    mp.text = matrix_text
    mp.font.name = "Consolas"
    mp.font.size = Pt(24)
    mp.font.color.rgb = Theme.title
    for lb, x0, y0 in [("A", 8.4, 2.1), ("D", 10.8, 2.1), ("B", 9.6, 4.4)]:
        add_rounded_box(s, x0, y0, 1.3, 0.8, lb, RGBColor(226, 232, 240), font_size=18)
    add_arrow(s, 9.55, 2.85, 1.6, 0.35, color=Theme.accent)
    add_arrow(s, 10.35, 3.65, -1.6, 0.35, color=Theme.accent2)
    add_arrow(s, 8.75, 3.65, 1.6, 0.35, color=Theme.warn)
    add_bullets(s, ["直覺：A 克 D、D 克 B、B 克 A，促進旋轉而非靜態均衡"], x=0.85, y=5.5, w=11.5, h=1.0, size=17)
    add_footer(s, slide_no)
    slide_no += 1

    # 9. Why oscillation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "為什麼會振盪？")
    add_rounded_box(s, 1.0, 2.2, 3.0, 1.0, "循環 payoff", RGBColor(254, 240, 138), font_size=18)
    add_rounded_box(s, 4.8, 2.2, 3.0, 1.0, "1-step delay", RGBColor(219, 234, 254), font_size=18)
    add_rounded_box(s, 8.6, 2.2, 3.0, 1.0, "有限抽樣噪音", RGBColor(209, 250, 229), font_size=18)
    add_arrow(s, 4.05, 2.5, 0.65, 0.4)
    add_arrow(s, 7.85, 2.5, 0.65, 0.4)
    add_bullets(
        s,
        [
            "reward(t) 依賴上一輪分布，形成延遲回饋路徑",
            "離散樣本 + 延遲常把理想封閉軌道變成近似振盪",
            "結果：Level 2 容易觀測，Level 3 較難穩定成立",
        ],
        x=1.0,
        y=4.0,
        w=11.3,
        h=2.0,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 10. Evolution rule
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "演化更新規則（Replicator-like）")
    eq = s.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(7.6), Inches(1.0))
    ep = eq.text_frame.paragraphs[0]
    ep.text = "w_s = exp(k * (u_s_bar - u_bar))，再 normalize 到 mean(weight)=1"
    ep.font.name = "Consolas"
    ep.font.size = Pt(20)
    ep.font.color.rgb = Theme.title
    add_bullets(
        s,
        [
            "selection_strength (k) 越大，反應越硬，可能放大振盪",
            "exp 形式確保權重為正，避免負權重問題",
            "對應程式層：evolution/replicator_dynamics.py",
        ],
        x=0.9,
        y=3.2,
        w=7.0,
        h=2.5,
        size=18,
    )
    add_rounded_box(s, 8.5, 2.2, 3.6, 1.1, "表現好：權重上升", RGBColor(220, 252, 231), font_size=17)
    add_rounded_box(s, 8.5, 3.7, 3.6, 1.1, "表現差：權重下降", RGBColor(254, 226, 226), font_size=17)
    add_footer(s, slide_no)
    slide_no += 1

    # 11. Economic model
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "經濟模型：資本流、資訊流、壓力流")
    add_highlight_box(s, 1.0, 2.1, 3.3, 1.2, "金幣\n行動權", kind="warning", font_size=18)
    add_highlight_box(s, 5.0, 2.1, 3.3, 1.2, "樣本\n學習能力", kind="info", font_size=18)
    add_highlight_box(s, 9.0, 2.1, 3.3, 1.2, "環境壓力\n選擇方向", kind="success", font_size=18)
    add_arrow(s, 4.35, 2.45, 0.55, 0.4)
    add_arrow(s, 8.35, 2.45, 0.55, 0.4)
    add_bullets(
        s,
        [
            "資源三句話：誰能動（行動權）、怎麼動（學習能力）、往哪動（選擇方向）",
            "金幣不直接進演化方程；樣本與壓力才決定動力學形狀",
        ],
        x=0.95,
        y=4.0,
        w=11.3,
        h=2.0,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 11.1 Failure cost precision slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "失敗成本的精確定義", "Failure = 經濟 + 學習 + 動力學 三層")
    add_highlight_box(s, 0.9, 2.1, 3.8, 1.2, "① 經濟層\n金幣消耗 → 行動能力下降", kind="danger", font_size=16)
    add_highlight_box(s, 4.9, 2.1, 3.8, 1.2, "② 學習層\n高價值樣本外流給對手", kind="danger", font_size=16)
    add_highlight_box(s, 8.9, 2.1, 3.8, 1.2, "③ 動力學層\n自身權重受 replicator 拉低", kind="danger", font_size=16)
    add_bullets(
        s,
        [
            "這個三層框架讓 gameplay、經濟、數學可在同一張圖對齊。",
            "也能直接對應實驗指標：資源、樣本、stage score 的變化。",
        ],
        x=1.0,
        y=4.2,
        w=11.2,
        h=1.4,
        size=17,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 12. Full integrated loop
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "完整整合循環（微觀 → 中觀 → 宏觀）")
    add_rounded_box(s, 1.0, 2.4, 2.2, 1.0, "玩家行為", RGBColor(219, 234, 254), font_size=17)
    add_rounded_box(s, 3.5, 2.4, 2.2, 1.0, "產生樣本", RGBColor(226, 232, 240), font_size=17)
    add_rounded_box(s, 6.0, 2.4, 2.2, 1.0, "小火龍分析", RGBColor(254, 240, 138), font_size=17)
    add_rounded_box(s, 8.5, 2.4, 2.2, 1.0, "生成地下城", RGBColor(209, 250, 229), font_size=17)
    add_rounded_box(s, 11.0, 2.4, 2.2, 1.0, "玩家適應", RGBColor(219, 234, 254), font_size=17)
    for x in [3.25, 5.75, 8.25, 10.75]:
        add_arrow(s, x, 2.65, 0.2, 0.35)
    add_connector(s, 12.0, 3.45, 2.0, 5.5)
    add_bullets(s, ["這是一個閉環 AI 生態系：行為會改變世界，世界也會反向塑造行為"], x=0.95, y=5.6, w=11.4, h=0.9, size=17)
    add_footer(s, slide_no)
    slide_no += 1

    # 13. Why non-convergent
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "為什麼系統不易收斂？", "可視為不收斂的三個必要條件")
    add_highlight_box(s, 1.2, 2.2, 3.4, 1.1, "1. 無單一最優解\n(cyclic payoff)", kind="warning", font_size=16)
    add_highlight_box(s, 5.0, 2.2, 3.4, 1.1, "2. 環境會反應\n(Meta GAI)", kind="info", font_size=16)
    add_highlight_box(s, 8.8, 2.2, 3.4, 1.1, "3. 玩家不一致\n(heterogeneity)", kind="success", font_size=16)
    add_bullets(
        s,
        [
            "三層共同作用形成 Dynamic Equilibrium，而非單點靜態均衡",
            "研究重點不是找到唯一最優策略，而是理解吸引子結構如何轉移",
        ],
        x=1.0,
        y=4.2,
        w=11.3,
        h=1.5,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 13.1 Cross-domain mapping
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "這個系統像什麼？", "跨領域對照：快速定位")
    add_highlight_box(s, 1.0, 2.0, 5.3, 1.0, "演化博弈：replicator dynamics", kind="info", font_size=17)
    add_highlight_box(s, 1.0, 3.2, 5.3, 1.0, "強化學習：policy-weight update", kind="success", font_size=17)
    add_highlight_box(s, 1.0, 4.4, 5.3, 1.0, "經濟學：市場競爭 + 資訊流", kind="warning", font_size=17)
    add_highlight_box(s, 7.0, 2.6, 5.3, 1.0, "控制系統：feedback loop", kind="danger", font_size=17)
    add_bullets(
        s,
        [
            "這頁用途：讓教授與跨域組員在 30 秒內知道它不是單一學門問題。",
        ],
        x=0.95,
        y=5.7,
        w=11.2,
        h=0.7,
        size=15,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 14. Metrics and levels
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "循環分級指標（研究契約）")
    add_bullets(
        s,
        [
            "Level 0：無循環訊號",
            "Level 1：有振幅（Amplitude）",
            "Level 2：振幅 + 主頻（Autocorr Dominant Frequency）",
            "Level 3：方向一致（Stage3 Phase Direction Consistency）",
        ],
        x=0.9,
        y=1.9,
        w=7.0,
        h=3.8,
        size=19,
    )
    add_rounded_box(s, 8.5, 2.0, 3.8, 1.0, "Stage 1: Amplitude", RGBColor(226, 232, 240), font_size=16)
    add_rounded_box(s, 8.5, 3.2, 3.8, 1.0, "Stage 2: Dominant Frequency", RGBColor(226, 232, 240), font_size=16)
    add_rounded_box(s, 8.5, 4.4, 3.8, 1.0, "Stage 3: Rotation Consistency", RGBColor(226, 232, 240), font_size=16)
    add_footer(s, slide_no)
    slide_no += 1

    # 15. Empirical results
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "目前實驗結果（對齊研發日誌）")
    add_bullets(
        s,
        [
            "大量條件下可穩定觀測 Level 2（振盪存在）",
            "Level 3 在 formal multi-seed protocol 幾乎為 0",
            "W2.1R、W3.1、W3.2、W3.3 全線 closure：未出現可承認 Level 3 seeds",
            "B4/B3/B5 進一步支持：局部機制有反應，但不足以改變 attractor class",
        ],
        x=0.9,
        y=1.9,
        w=7.6,
        h=4.2,
        size=17,
    )
    add_rounded_box(s, 9.2, 2.1, 2.8, 1.0, "Level 2\n大量存在", RGBColor(220, 252, 231), font_size=17)
    add_rounded_box(s, 9.2, 3.4, 2.8, 1.0, "Level 3\n幾乎 0", RGBColor(254, 226, 226), font_size=17)
    add_rounded_box(s, 9.2, 4.7, 2.8, 1.0, "結論\nPlateau 穩固", RGBColor(254, 240, 138), font_size=17)
    add_footer(s, slide_no)
    slide_no += 1

    # 15.1 L2 vs L3 visual intuition
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "Level 2 vs Level 3：直覺圖像")
    add_highlight_box(s, 0.9, 1.9, 5.7, 0.9, "Level 2：有振盪，但方向不穩定（亂晃）", kind="warning", font_size=16)
    add_highlight_box(s, 6.8, 1.9, 5.7, 0.9, "Level 3：方向一致的穩定旋轉（順/逆時針）", kind="success", font_size=16)
    circle_l2 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.5), Inches(3.0), Inches(3.8), Inches(2.4))
    circle_l2.fill.solid()
    circle_l2.fill.fore_color.rgb = RGBColor(254, 243, 199)
    circle_l2.line.color.rgb = RGBColor(245, 158, 11)
    circle_l3 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.0), Inches(3.0), Inches(3.8), Inches(2.4))
    circle_l3.fill.solid()
    circle_l3.fill.fore_color.rgb = RGBColor(220, 252, 231)
    circle_l3.line.color.rgb = RGBColor(16, 185, 129)
    add_arrow(s, 2.1, 3.7, 1.0, 0.35)
    add_arrow(s, 3.1, 4.2, -0.7, 0.35)
    add_arrow(s, 4.0, 3.7, 0.7, 0.35)
    add_arrow(s, 8.8, 3.9, 2.0, 0.35)
    add_bullets(
        s,
        [
            "研究核心問題：我們目前大量到達 Level 2，但尚未把旋轉方向長窗鎖定到 Level 3。",
        ],
        x=0.95,
        y=5.8,
        w=11.4,
        h=0.7,
        size=15,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 16. Bottlenecks
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "目前瓶頸：為什麼卡在 Level 2 Plateau")
    add_rounded_box(s, 1.1, 2.3, 3.5, 1.1, "記憶太短\n(主要是 1-step lag)", RGBColor(254, 226, 226), font_size=16)
    add_rounded_box(s, 4.9, 2.3, 3.5, 1.1, "更新仍偏均化\n(shared update pressure)", RGBColor(254, 226, 226), font_size=16)
    add_rounded_box(s, 8.7, 2.3, 3.5, 1.1, "相位方向難鎖定\n(score 無法長窗穩定)", RGBColor(254, 226, 226), font_size=16)
    add_bullets(
        s,
        [
            "已排除：單純調 a/b、固定 commitment、hysteresis、pulse 等修補路線",
            "現階段最可能是結構性問題：需要能改寫 basin topology 的新機制",
        ],
        x=1.0,
        y=4.3,
        w=11.2,
        h=1.6,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 17. Next H-series
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "下一步研發主線（H-series）")
    add_rounded_box(s, 1.0, 2.2, 3.7, 1.4, "H1 記憶機制\n多輪歷史與慣性", RGBColor(219, 234, 254), font_size=17)
    add_rounded_box(s, 4.9, 2.2, 3.7, 1.4, "H2 非線性世界\nthreshold/regime 切換", RGBColor(209, 250, 229), font_size=17)
    add_rounded_box(s, 8.8, 2.2, 3.7, 1.4, "H3 異質玩家\n不同學習速率/族群", RGBColor(254, 240, 138), font_size=17)
    add_bullets(
        s,
        [
            "目標：不再做局部 tuning，而是改變相位幾何與吸引子結構",
            "驗證標準：multi-seed Level 3 + stage3 score + env_gamma 一致提升",
        ],
        x=1.0,
        y=4.4,
        w=11.2,
        h=1.4,
        size=18,
    )
    add_footer(s, slide_no)
    slide_no += 1

    # 18. Work blocks
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "可拆分的大工作區塊（建議分工）")
    blocks = [
        ("Block 1\n核心模擬引擎", "core/ players/ dungeon/ evolution/", RGBColor(219, 234, 254)),
        ("Block 2\n實驗與掃參", "simulation/ scripts/ outputs/", RGBColor(209, 250, 229)),
        ("Block 3\n分析指標", "analysis/ cycle metrics/ plotting", RGBColor(254, 240, 138)),
        ("Block 4\n機制設計", "H-series / policy / new dynamics", RGBColor(254, 226, 226)),
        ("Block 5\n敘事產品化", "docs/ story/ UI prototype", RGBColor(226, 232, 240)),
    ]
    y = 1.9
    for title, desc, color in blocks:
        add_rounded_box(s, 0.9, y, 5.8, 0.9, title, color, font_size=15)
        add_rounded_box(s, 7.0, y, 5.4, 0.9, desc, RGBColor(241, 245, 249), font_size=14, bold=False)
        y += 1.05
    add_footer(s, slide_no)
    slide_no += 1

    # 19. SDD invariants (auto split demo for dense text)
    sdd_lines = [
        "先 Spec 後改碼：變更行為契約必先更新 SDD.md",
        "分層不變條件：analysis 不反向依賴 simulation；evolution 不做 I/O",
        "所有可重現命令在 venv：./venv/bin/python -m ...",
        "輸出契約固定：timeseries CSV 欄位與 burn-in/tail 視窗語意一致",
        "研究結論以 研發日誌.md 與 outputs/ 產物交叉驗證",
    ]
    slide_no = add_auto_paged_bullets(
        prs,
        slide_no,
        "開發與研究契約（SDD 必守）",
        "高密度內容示範：超過 4 行自動分頁",
        sdd_lines,
        x=0.9,
        y=1.8,
        w=11.4,
        h=4.3,
        size=17,
        max_lines=4,
    )

    # 20. Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide_bg(s)
    add_title_box(s, "最終總結", "這是一個永遠不會輕易收斂的 AI 生態系")
    add_bullets(
        s,
        [
            "玩家在學習，世界也在學習",
            "地下城不是靜態關卡，而是人格互動的動態投影",
            "目前已證明 plateau 穩固；下一步是結構性突破而非局部調參",
        ],
        x=1.0,
        y=2.2,
        w=11.2,
        h=2.7,
        size=24,
    )
    add_rounded_box(s, 3.0, 5.1, 7.4, 1.2, "Q&A | 可接續：15 分鐘版、30 分鐘版、分工版", RGBColor(209, 250, 229), font_size=22)
    add_footer(s, slide_no)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    out = build_deck()
    print(f"Generated: {out}")
