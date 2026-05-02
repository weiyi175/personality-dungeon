from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


RESEARCH_OUT = Path("docs/personality_dungeon_research_30min_with_notes.pptx")
TEAM_OUT = Path("docs/personality_dungeon_team_workplan.pptx")
TEAM_MD = Path("docs/team_block_workplan.md")


class C:
    BG = RGBColor(248, 250, 252)
    TITLE = RGBColor(15, 23, 42)
    BODY = RGBColor(30, 41, 59)
    ACCENT = RGBColor(14, 116, 144)
    BOX = RGBColor(226, 232, 240)
    GOOD = RGBColor(220, 252, 231)
    WARN = RGBColor(254, 240, 138)
    BAD = RGBColor(254, 226, 226)
    BLUE = RGBColor(219, 234, 254)
    GREEN = RGBColor(209, 250, 229)


@dataclass
class SlideSpec:
    title: str
    subtitle: str
    bullets: list[str]
    note: str


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C.BG


def title(slide, t: str, st: str = ""):
    tbox = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.9))
    p = tbox.text_frame.paragraphs[0]
    p.text = t
    p.font.name = "Microsoft JhengHei"
    p.font.bold = True
    p.font.size = Pt(33)
    p.font.color.rgb = C.TITLE

    if st:
        sbox = slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.6), Inches(0.6))
        sp = sbox.text_frame.paragraphs[0]
        sp.text = st
        sp.font.name = "Microsoft JhengHei"
        sp.font.size = Pt(16)
        sp.font.color.rgb = C.ACCENT


def bullets(slide, lines: list[str], x=0.85, y=1.75, w=7.0, h=4.8, size=19):
    def _render_lines(target_x, target_y, target_w, target_h, target_lines, target_size):
        box = slide.shapes.add_textbox(Inches(target_x), Inches(target_y), Inches(target_w), Inches(target_h))
        tf = box.text_frame
        tf.clear()
        for i, line in enumerate(target_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Microsoft JhengHei"
            p.font.size = Pt(target_size)
            p.font.color.rgb = C.BODY
            p.space_after = Pt(7)

    if len(lines) <= 4:
        _render_lines(x, y, w, h, lines, size)
        return

    # Auto split when content is dense: two columns, smaller font.
    mid = (len(lines) + 1) // 2
    left = lines[:mid]
    right = lines[mid:]
    gap = 0.35
    col_w = (w - gap) / 2
    _render_lines(x, y, col_w, h, left, max(15, size - 2))
    _render_lines(x + col_w + gap, y, col_w, h, right, max(15, size - 2))


def rbox(slide, x, y, w, h, txt, fill, fs=16, bold=True):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(148, 163, 184)
    p = shp.text_frame.paragraphs[0]
    p.text = txt
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft JhengHei"
    p.font.size = Pt(fs)
    p.font.bold = bold
    p.font.color.rgb = C.TITLE
    return shp


def arrow(slide, x, y, w=0.6, h=0.35):
    a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor(100, 116, 139)
    a.line.color.rgb = RGBColor(100, 116, 139)


def footer(slide, text: str):
    f = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.8), Inches(0.25))
    p = f.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def add_note(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def make_research_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    specs: list[SlideSpec] = [
        SlideSpec(
            title="個性地下城：30 分鐘研究版",
            subtitle="人格 × 演化 × 系統動力學（含講稿與 Q&A）",
            bullets=[
                "目標：讓新組員在 30 分鐘內掌握研究問題、模型、結果、下一步",
                "資料基礎：SDD 契約 + 研發日誌正式 closure 結論",
                "核心狀態：Level 2 plateau 穩固，Level 3 emergence 未正式出現",
            ],
            note=(
                "開場 1 分鐘：\n"
                "1) 先說這不是做 RPG 成品，而是研究型動力系統。\n"
                "2) 今日重點是可重現證據，不是單次漂亮曲線。\n"
                "3) 先給結論：我們已系統化排除多條直覺路線。"
            ),
        ),
        SlideSpec(
            title="一句話框架",
            subtitle="玩家在學習，世界也在學習",
            bullets=[
                "玩家：個體演化單元（策略、效用、權重）",
                "地下城：人格測試場（行為回饋）",
                "小火龍：世界壓力控制器（反收斂）",
            ],
            note="這頁講 1 分鐘，目的只做共同語言。不要陷入細節。",
        ),
        SlideSpec(
            title="研究問題（Research Question）",
            subtitle="從振盪走到穩定方向循環是否可能？",
            bullets=[
                "RQ1：為什麼 Level 2 很常見？",
                "RQ2：為什麼 Level 3 幾乎不出現？",
                "RQ3：哪些機制改動能改寫 attractor class？",
            ],
            note="強調我們關心的是吸引子結構，不是單次分數。",
        ),
        SlideSpec(
            title="玩家機制與迴圈",
            subtitle="選擇策略 → 地下城回饋 → 演化更新 → 下一輪",
            bullets=[
                "每輪輸出 timeseries：round, p_*, w_*, avg_reward, avg_utility",
                "玩家採樣來自 strategy weights，回饋來自 dungeon payoff",
                "演化更新採 exp-weight replicator-like",
            ],
            note="講 2 分鐘，補充：simulation 層負責 I/O 與 CSV 契約。",
        ),
        SlideSpec(
            title="數學模型 I：狀態與延遲",
            subtitle="x(t) = (x_A, x_D, x_B), sum(x)=1",
            bullets=[
                "策略集合：aggressive / defensive / balanced",
                "第 t 輪 reward 主要使用上一輪 popularity（1-step lag）",
                "延遲回饋提高振盪可能，也提高噪音敏感性",
            ],
            note="提醒聽眾：這是理論與程式順序對齊的關鍵。",
        ),
        SlideSpec(
            title="數學模型 II：Payoff 矩陣",
            subtitle="U = A x（RPS-like 旋轉結構）",
            bullets=[
                "A = [[0, a, -b], [-b, 0, a], [a, -b, 0]]",
                "直覺：A 克 D、D 克 B、B 克 A",
                "理論上支持循環，但離散與噪音會改變軌道品質",
            ],
            note="此頁 2 分鐘，建議黑板手畫 A→D→B→A。",
        ),
        SlideSpec(
            title="數學如何對應玩家行為",
            subtitle="x(t) → U=Ax → reward → 權重更新 → 下一輪分布",
            bullets=[
                "x(t)：玩家群體策略比例（狀態）",
                "U=Ax：地下城對該狀態給出的回饋",
                "reward 高：該策略權重上升，玩家比例增加",
                "分布變了：下一輪 reward 結構也會改變",
                "結論：玩家在玩遊戲，同時在解延遲動態方程",
            ],
            note="這頁是關鍵橋樑，專門把『聽懂兩邊卻連不起來』的問題解掉。",
        ),
        SlideSpec(
            title="演化更新",
            subtitle="w_s = exp(k*(u_s_bar - u_bar))，並做 normalize",
            bullets=[
                "k 越高：放大差異，也可能放大 overshoot",
                "exp 確保權重始終為正值",
                "mean(weight)=1 用來避免尺度漂移",
            ],
            note="補充工程合理性：避免線性更新出現負權重。",
        ),
        SlideSpec(
            title="經濟模型：三種流",
            subtitle="金幣流 + 樣本流 + 壓力流",
            bullets=[
                "金幣 = 行動權（誰能動）",
                "樣本 = 學習能力（怎麼動）",
                "壓力 = 選擇方向（往哪動）",
            ],
            note="強調：金幣不直接進演化方程，樣本才是長期核心。",
        ),
        SlideSpec(
            title="失敗成本三層（精確版）",
            subtitle="經濟層 + 學習層 + 動力學層",
            bullets=[
                "經濟層：金幣消耗，後續行動能力下降",
                "學習層：高價值樣本外流到對手",
                "動力學層：自身策略權重被 replicator 拉低",
                "研究價值：同時連到經濟模型、數學模型、玩家行為",
            ],
            note="這頁可直接回應『失敗到底損失什麼』。",
        ),
        SlideSpec(
            title="圖例（Legend）",
            subtitle="顏色語意統一",
            bullets=[
                "藍：玩家與狀態",
                "綠：地下城 / 正向回饋",
                "黃：小火龍 / 控制器",
                "紅：風險與失敗成本",
            ],
            note="先講這頁可降低後面跨頁理解成本。",
        ),
        SlideSpec(
            title="系統分層（SDD Invariants）",
            subtitle="為何我們能信任實驗結果",
            bullets=[
                "analysis 不反向 import simulation",
                "evolution 不做 I/O、不依賴 plotting",
                "simulation 管資料契約與輸出欄位",
            ],
            note="這頁 1 分鐘，目的是建立可重現性與責任邊界。",
        ),
        SlideSpec(
            title="循環分級指標",
            subtitle="Level 0/1/2/3 的操作化定義",
            bullets=[
                "Stage1：Amplitude（峰谷振幅）",
                "Stage2：Dominant Frequency（自相關主頻）",
                "Stage3：Phase Direction Consistency（方向一致）",
            ],
            note="提醒：Stage2 具 min_overlap/min_cycles guard，避免假陽性。",
        ),
        SlideSpec(
            title="Level 2 vs Level 3：直覺圖像",
            subtitle="亂晃振盪 vs 穩定旋轉",
            bullets=[
                "Level 2：振幅與主頻存在，但方向一致性不足",
                "Level 3：存在穩定順/逆時針旋轉，方向一致",
                "核心研究問題：如何把 Level 2 轉成可穩健的 Level 3",
            ],
            note="這頁建議多停 30 秒，幫新手形成圖像。",
        ),
        SlideSpec(
            title="實驗總覽：我們觀測到什麼",
            subtitle="Level 2 常見，Level 3 稀有",
            bullets=[
                "多 seed / 多格點結果：Level 2 可穩定出現",
                "formal protocol 下 Level 3 seeds 長期接近 0",
                "結論：系統有振盪，但方向性難以長窗鎖定",
            ],
            note="這頁要用『有現象但不突破』敘事，保持科學中立。",
        ),
        SlideSpec(
            title="W2.1R closure（生存修復線）",
            subtitle="修復了生存，不等於產生 emergence",
            bullets=[
                "life 4..6 可跑滿，但 tail Level3 seed 仍為 0",
                "personality drift 朝保守存活方向集中",
                "正式 decision：close_w2_1",
            ],
            note="Q&A 常問：是不是只要再長跑就會有？答：formal gate 已涵蓋。",
        ),
        SlideSpec(
            title="W3.1~W3.3 closure（leader policy 線）",
            subtitle="靜態 commitment / hysteresis / pulse 全部結案",
            bullets=[
                "policy 有啟動，不是 no-op",
                "但 seed-level level counts 不變",
                "正式結論：仍停在同一 Level 2 plateau",
            ],
            note="關鍵答法：『有局部 uplift，不代表 attractor class 改變。』",
        ),
        SlideSpec(
            title="B4/B3/B5 closure（sampled-side 主力線）",
            subtitle="局部差異可注入，但不足以打開 basin",
            bullets=[
                "B4：state-dependent k 訊號太弱",
                "B3：strata 分離存在，但 shared update 壓回 plateau",
                "B5：deterministic 通過，sampled 仍無 Level 3",
            ],
            note="這頁是研究價值核心：排除了『只是沒打進去』的質疑。",
        ),
        SlideSpec(
            title="目前最堅實結論",
            subtitle="Attractor Robustness > 參數微調",
            bullets=[
                "多條路線皆有機制痕跡，但皆未改變 seed-level attractor",
                "問題已從 tuning 問題升級為結構問題",
                "後續應聚焦可改寫 basin topology 的機制",
            ],
            note="此頁 2 分鐘，建議做口語收斂：不是失敗，是結構性發現。",
        ),
        SlideSpec(
            title="為什麼不收斂：三個必要條件",
            subtitle="像 theorem 一樣讀",
            bullets=[
                "1) 無單一最優解：cyclic payoff",
                "2) 環境會回應：Meta GAI 介入",
                "3) 玩家不一致：heterogeneity 存在",
            ],
            note="這頁用來提升專業度與論證力度。",
        ),
        SlideSpec(
            title="這個系統像什麼？",
            subtitle="跨領域定位",
            bullets=[
                "演化博弈：replicator dynamics",
                "強化學習：policy update",
                "經濟學：市場競爭 + 資訊流",
                "控制系統：feedback loop",
            ],
            note="跨領域聽眾會在這頁快速找到定位。",
        ),
        SlideSpec(
            title="下一步研究設計（H 系列）",
            subtitle="不是微調，而是機制層改寫",
            bullets=[
                "H1：更長記憶核與歷史慣性",
                "H2：nonlinear threshold/regime dynamics",
                "H3：異質玩家學習率與亞群耦合",
            ],
            note="說明驗收標準：multi-seed Level3 + score + env_gamma。",
        ),
        SlideSpec(
            title="工作分塊（研究工程化）",
            subtitle="核心模擬 / 掃參實驗 / 分析指標 / 機制設計 / 敘事整合",
            bullets=[
                "Block 1：Simulation Core（不可亂改）",
                "Block 2：Experiment System（研究主力）",
                "Block 3：Analysis（論文核心）",
                "Block 4：Mechanism Design（突破關鍵）",
                "Block 5：Narrative & Productization（外溝通）",
            ],
            note="過渡到第二份組員分工版簡報。",
        ),
        SlideSpec(
            title="Q&A 備答總表（研究版）",
            subtitle="先備 6 題，避免現場失焦",
            bullets=[
                "Q1：是不是門檻太嚴？A：已做多閾值與 guard calibration",
                "Q2：是不是 seed 不夠多？A：formal protocol 已做 multi-seed",
                "Q3：是不是跑不夠久？A：長窗與 closure gate 已覆蓋",
                "Q4：是否可用 RL 直接解？A：可，但要先定義新 family 與對照",
                "Q5：為何不繼續微調 W3？A：closure 規則已鎖，避免 p-hacking",
                "Q6：下一步成功判準？A：seed-level Level3 + 指標一致 uplift",
            ],
            note="建議把這頁當備用，不必逐條念完。",
        ),
        SlideSpec(
            title="結尾",
            subtitle="這是一個可重現、可防守、可擴展的 AI 生態研究平台",
            bullets=[
                "已完成：現象建立 + 負結果收斂 + closure 敘事",
                "正在做：新 family 方向設計與分工落地",
                "需要組員：機制、分析、實驗、敘事四線並行",
            ],
            note="收尾 1 分鐘：請組員加入具體 block 任務。",
        ),
    ]

    for idx, sp in enumerate(specs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide)
        title(slide, sp.title, sp.subtitle)
        bullets(slide, sp.bullets)
        slide_title = sp.title

        # Right-side visual aids by section
        if slide_title in ("一句話框架", "研究問題（Research Question）"):
            rbox(slide, 8.2, 2.0, 4.1, 1.0, "玩家", C.BLUE, fs=20)
            rbox(slide, 8.2, 3.2, 4.1, 1.0, "地下城", C.GREEN, fs=20)
            rbox(slide, 8.2, 4.4, 4.1, 1.0, "小火龍", C.WARN, fs=20)
        elif slide_title == "玩家機制與迴圈":
            xs = [8.0, 9.1, 10.2, 11.3]
            names = ["選策略", "算回饋", "更新權重", "下一輪"]
            for x0, n in zip(xs, names):
                rbox(slide, x0, 2.8, 1.0, 1.0, n, C.BOX, fs=12)
            arrow(slide, 9.0, 3.1, 0.1, 0.25)
            arrow(slide, 10.1, 3.1, 0.1, 0.25)
            arrow(slide, 11.2, 3.1, 0.1, 0.25)
        elif slide_title in ("數學模型 I：狀態與延遲", "數學模型 II：Payoff 矩陣", "數學如何對應玩家行為", "演化更新"):
            rbox(slide, 8.1, 2.1, 4.1, 1.1, "數學重點", C.BOX, fs=16)
            if slide_title == "數學模型 II：Payoff 矩陣":
                m = slide.shapes.add_textbox(Inches(8.2), Inches(3.35), Inches(4.0), Inches(2.1))
                mp = m.text_frame.paragraphs[0]
                mp.text = "A = [0  a -b]\n    [-b 0  a]\n    [a -b 0]"
                mp.font.name = "Consolas"
                mp.font.size = Pt(22)
                mp.font.color.rgb = C.TITLE
            elif slide_title == "數學如何對應玩家行為":
                rbox(slide, 8.1, 3.4, 4.1, 0.9, "x(t) -> U=Ax", C.WARN, fs=16)
                rbox(slide, 8.1, 4.5, 4.1, 0.9, "reward -> weight", C.GOOD, fs=16)
            else:
                rbox(slide, 8.1, 3.5, 4.1, 1.8, "圖示區", C.BLUE, fs=14)
        elif slide_title in ("循環分級指標", "Level 2 vs Level 3：直覺圖像"):
            rbox(slide, 8.2, 2.2, 4.0, 1.0, "Stage 1", C.BOX, fs=18)
            rbox(slide, 8.2, 3.4, 4.0, 1.0, "Stage 2", C.BOX, fs=18)
            rbox(slide, 8.2, 4.6, 4.0, 1.0, "Stage 3", C.BOX, fs=18)
            if slide_title == "Level 2 vs Level 3：直覺圖像":
                rbox(slide, 8.2, 5.8, 4.0, 0.7, "L2: drift / L3: stable rotation", C.WARN, fs=12)
        elif slide_title in ("實驗總覽：我們觀測到什麼", "W2.1R closure（生存修復線）", "W3.1~W3.3 closure（leader policy 線）", "B4/B3/B5 closure（sampled-side 主力線）", "目前最堅實結論"):
            rbox(slide, 8.2, 2.1, 4.0, 1.2, "Level 2: common", C.GOOD, fs=18)
            rbox(slide, 8.2, 3.6, 4.0, 1.2, "Level 3: rare", C.BAD, fs=18)
            rbox(slide, 8.2, 5.1, 4.0, 1.0, "Plateau robust", C.WARN, fs=17)
        elif slide_title == "下一步研究設計（H 系列）":
            rbox(slide, 8.2, 2.2, 4.1, 1.1, "H1", C.BLUE, fs=19)
            rbox(slide, 8.2, 3.6, 4.1, 1.1, "H2", C.GREEN, fs=19)
            rbox(slide, 8.2, 5.0, 4.1, 1.1, "H3", C.WARN, fs=19)
        elif slide_title == "工作分塊（研究工程化）":
            y = 1.9
            for bname in ["Block 1", "Block 2", "Block 3", "Block 4", "Block 5"]:
                rbox(slide, 8.0, y, 4.3, 0.85, bname, C.BOX, fs=14)
                y += 0.95
        elif slide_title == "圖例（Legend）":
            rbox(slide, 8.2, 2.0, 4.1, 0.9, "藍=玩家", C.BLUE, fs=14)
            rbox(slide, 8.2, 3.1, 4.1, 0.9, "綠=地下城", C.GREEN, fs=14)
            rbox(slide, 8.2, 4.2, 4.1, 0.9, "黃=小火龍", C.WARN, fs=14)
            rbox(slide, 8.2, 5.3, 4.1, 0.9, "紅=風險/失敗", C.BAD, fs=14)
        elif slide_title == "這個系統像什麼？":
            rbox(slide, 8.2, 2.2, 4.1, 1.0, "Game Theory", C.BLUE, fs=16)
            rbox(slide, 8.2, 3.5, 4.1, 1.0, "Reinforcement Learning", C.GREEN, fs=16)
            rbox(slide, 8.2, 4.8, 4.1, 1.0, "Control System", C.WARN, fs=16)

        footer(slide, f"Research Deck 30 min | p.{idx}")
        add_note(slide, sp.note)

    RESEARCH_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(RESEARCH_OUT))
    return RESEARCH_OUT


def make_team_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "組員分工版：Block 任務地圖", "每個 Block 對應任務、輸入/輸出、預估工時")
    bullets(
        s,
        [
            "用途：快速 onboarding + 任務切分 + 排程對齊",
            "工時單位：人時（ideal），不含等待 queue 時間",
            "輸入/輸出以可驗證檔案或指標為準",
        ],
        x=0.85,
        y=1.85,
        w=7.2,
        h=3.2,
        size=18,
    )
    rbox(s, 8.3, 2.0, 3.9, 1.0, "Block 1-5", C.WARN, fs=20)
    rbox(s, 8.3, 3.3, 3.9, 1.0, "Task I/O", C.BLUE, fs=20)
    rbox(s, 8.3, 4.6, 3.9, 1.0, "Est. Hours", C.GREEN, fs=20)
    footer(s, "Team Workplan Deck | p.1")
    add_note(s, "開場 30 秒：這份是協作用，不是理論介紹。")

    block_rows = [
        ("Block 1 核心模擬引擎", "修正核心 loop、事件掛點、保持分層 invariants", "SDD.md 契約、core/players/dungeon/evolution", "可跑 simulation + 無契約破壞 PR", "24-40h"),
        ("Block 2 實驗與掃參", "seed stability、grid/refine、formal gate", "baseline summary、protocol 參數", "outputs/*.csv + summary json/tsv", "30-50h"),
        ("Block 3 分析指標", "cycle metrics、stage3 guards、可視化", "timeseries CSV、analysis 契約", "診斷報表、圖表、pass/fail 判定", "24-36h"),
        ("Block 4 機制設計", "H1/H2/H3 新 family 最小實作", "closure 結論、失敗模式清單", "新機制實驗結果 + ablation", "40-64h"),
        ("Block 5 敘事與文件", "runbook、週報、對外簡報與Q&A", "研發日誌、outputs、分析圖", "可審核文件與投影片", "16-28h"),
    ]

    # Slides 2-6
    for i, (name, task, inp, out, hrs) in enumerate(block_rows, start=2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg(s)
        title(s, name, "任務 / 輸入 / 輸出 / 工時")
        rbox(s, 0.9, 1.9, 2.0, 0.8, "任務", C.BLUE, fs=15)
        rbox(s, 3.1, 1.9, 4.3, 0.8, "輸入", C.GREEN, fs=15)
        rbox(s, 7.6, 1.9, 3.6, 0.8, "輸出", C.WARN, fs=15)
        rbox(s, 11.4, 1.9, 1.0, 0.8, "工時", C.BOX, fs=15)

        rbox(s, 0.9, 2.9, 2.0, 2.8, task, RGBColor(241, 245, 249), fs=12, bold=False)
        rbox(s, 3.1, 2.9, 4.3, 2.8, inp, RGBColor(241, 245, 249), fs=12, bold=False)
        rbox(s, 7.6, 2.9, 3.6, 2.8, out, RGBColor(241, 245, 249), fs=12, bold=False)
        rbox(s, 11.4, 2.9, 1.0, 2.8, hrs, RGBColor(241, 245, 249), fs=12, bold=True)

        bullets(
            s,
            [
                "Definition of Done：輸出可重跑、指標可比對、紀錄可追溯",
                "交接格式：一頁摘要 + 指令 + 主要輸出檔清單",
            ],
            x=0.95,
            y=5.95,
            w=11.3,
            h=0.9,
            size=14,
        )
        footer(s, f"Team Workplan Deck | p.{i}")
        add_note(s, f"這頁講 1.5 分鐘，強調 {name} 的邊界與交付。")

    # Slide 7: dependency
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Block 依賴關係", "先後順序與同步點")
    rbox(s, 1.0, 2.3, 2.2, 0.9, "Block 1", C.BLUE, fs=16)
    rbox(s, 3.8, 2.3, 2.2, 0.9, "Block 2", C.GREEN, fs=16)
    rbox(s, 6.6, 2.3, 2.2, 0.9, "Block 3", C.WARN, fs=16)
    rbox(s, 9.4, 2.3, 2.2, 0.9, "Block 4", C.BAD, fs=16)
    rbox(s, 5.2, 4.1, 2.8, 0.9, "Block 5", C.BOX, fs=16)
    arrow(s, 3.25, 2.55)
    arrow(s, 6.05, 2.55)
    arrow(s, 8.85, 2.55)
    bullets(
        s,
        [
            "推薦節奏：1+2+3 並行基線，4 走候選機制，5 做週期化輸出整合",
            "同步點：每週一次 gate（指標與輸出檔一致性）",
        ],
        x=0.95,
        y=5.3,
        w=11.2,
        h=1.2,
        size=17,
    )
    footer(s, "Team Workplan Deck | p.7")
    add_note(s, "這頁用來排人力：誰是 owner、誰是 reviewer。")

    # Slide 8: 2-week sprint
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "兩週衝刺建議排程", "Week 1 建基線、Week 2 做機制驗證")
    bullets(
        s,
        [
            "Week 1：Block1/2/3 建可重跑基線 + Block5 做中期彙整",
            "Week 2：Block4 實作 H1/H2/H3 最小版 + Block2/3 回測",
            "收斂條件：至少一條機制線通過預定 gate 或形成高品質負結果",
        ],
        x=0.9,
        y=2.0,
        w=8.0,
        h=3.2,
        size=18,
    )
    rbox(s, 9.1, 2.2, 3.2, 1.0, "Week 1\nBaseline", C.BLUE, fs=17)
    rbox(s, 9.1, 3.6, 3.2, 1.0, "Week 2\nMechanism", C.GREEN, fs=17)
    rbox(s, 9.1, 5.0, 3.2, 1.0, "Gate Review", C.WARN, fs=17)
    footer(s, "Team Workplan Deck | p.8")
    add_note(s, "可依組員數把工時切成 owner + support。")

    # Slide 9: risk matrix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "風險矩陣與備援", "避免無效迭代")
    bullets(
        s,
        [
            "風險 1：只看單 seed 好結果 -> 備援：強制 multi-seed gate",
            "風險 2：改碼破壞契約 -> 備援：SDD 檢核 + regression",
            "風險 3：局部調參循環 -> 備援：closure rule 先鎖定",
            "風險 4：輸出不可重現 -> 備援：固定 venv 指令與 protocol log",
        ],
        x=0.9,
        y=1.9,
        w=11.3,
        h=3.8,
        size=18,
    )
    footer(s, "Team Workplan Deck | p.9")
    add_note(s, "重點是流程風險，不是模型風險。")

    # Slide 10: close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "分工版結論", "每個 Block 都有可驗證輸出，才能加速研究")
    bullets(
        s,
        [
            "先對齊契約，再開工；先有 gate，再談突破",
            "成功不是只看 Level 3，而是可防守的證據品質",
            "下一步：指派 owner、凍結兩週任務、開始執行",
        ],
        x=1.0,
        y=2.3,
        w=10.8,
        h=2.6,
        size=22,
    )
    rbox(s, 3.2, 5.3, 7.0, 1.0, "Action: 今天完成 Owner 指派與週節點", C.GREEN, fs=20)
    footer(s, "Team Workplan Deck | p.10")
    add_note(s, "收尾：直接進 owner 指派。")

    TEAM_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(TEAM_OUT))
    return TEAM_OUT


def write_team_markdown() -> Path:
    lines = [
        "# 組員分工工作表（Block 對應）",
        "",
        "| Block | 主要任務 | 主要輸入 | 主要輸出 | 預估工時 |",
        "|---|---|---|---|---|",
        "| Block 1 核心模擬引擎 | 修正核心 loop、事件掛點、維持分層 invariants | SDD.md、core/ players/ dungeon/ evolution/ | 可重跑 simulation、無契約破壞變更 | 24-40h |",
        "| Block 2 實驗與掃參 | seed stability、grid/refine、formal gate | baseline summary、protocol 參數 | outputs CSV/TSV/JSON 與 gate 結果 | 30-50h |",
        "| Block 3 分析指標 | cycle metrics、stage3 guards、可視化 | timeseries CSV、analysis 規範 | 分析圖、判定報表、結論摘要 | 24-36h |",
        "| Block 4 機制設計 | H1/H2/H3 新 family 最小實作與 ablation | closure 結論、失敗模式 | 新機制報告與可重現實驗 | 40-64h |",
        "| Block 5 敘事文件 | runbook、週報、簡報、Q&A | 研發日誌、outputs、分析圖 | 對內外一致敘事文件 | 16-28h |",
        "",
        "## 兩週建議排程",
        "",
        "1. Week 1：Block1/2/3 建 baseline 與回歸檢核，Block5 同步整理。",
        "2. Week 2：Block4 跑機制最小版，Block2/3 依 gate 做回測。",
        "3. 週末 gate：確認輸出可重跑、指標可比較、結論可防守。",
    ]
    TEAM_MD.parent.mkdir(parents=True, exist_ok=True)
    TEAM_MD.write_text("\n".join(lines), encoding="utf-8")
    return TEAM_MD


if __name__ == "__main__":
    r = make_research_deck()
    t = make_team_deck()
    m = write_team_markdown()
    print(f"Generated: {r}")
    print(f"Generated: {t}")
    print(f"Generated: {m}")
